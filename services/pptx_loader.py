from abstract_file_loader import AbstractFileLoader
import os
from pptx import Presentation

class PPTXLoader(AbstractFileLoader):
    
    def __init__(self):
        super().__init__()
    
    def load(self, file_path, file_name):
        """
        Load the content from the specified PPTX file path.
        :param file_path: The path to the PPTX file to check.
        :param file_name: The name of the PPTX file to check.
        :return: A dictionary containing the file's metadata and content.
        """
        full_path = os.path.join(file_path, file_name)
        presentation = Presentation(full_path)
        
        content = []
        for slide in presentation.slides:
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content.append(shape.text.strip())
            content.append(" ".join(slide_content))
        
        doc = {
            "name": file_name,
            "path": full_path,
            "extension": "pptx",
            "content": "\n".join(content).strip().replace("\n", " "),
            "size_bytes": os.path.getsize(full_path)
        }
        
        return doc
    
    def exists(self, file_path, file_name):
        """
        Check if the specified PPTX file exists.

        :param file_path: The path to the PPTX file to check.
        :param file_name: The name of the PPTX file to check.
        :return: True if the file exists, False otherwise.
        """
        return os.path.exists(os.path.join(file_path, file_name)) and os.path.isfile(os.path.join(file_path, file_name))